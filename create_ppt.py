from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Business Dashboard Design"
subtitle.text = "Interactive Dashboard for Business Stakeholders\n\nCreated using Streamlit\nDataset: Superstore Sales Data"

# Slide 2: Objective
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Objective"
tf = body_shape.text_frame
tf.text = "Design an interactive dashboard for business stakeholders to inform data-driven decisions."

p = tf.add_paragraph()
p.text = "• Learn how to create dashboards that inform business decisions"
p.level = 1

p = tf.add_paragraph()
p.text = "• Use Power BI / Tableau tools (Streamlit alternative used)"
p.level = 1

# Slide 3: Dataset Overview
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Dataset Overview"
tf = body_shape.text_frame
tf.text = "Superstore Sales Dataset - Sample sales and financial data"

p = tf.add_paragraph()
p.text = "• Contains sales transactions with customer, product, and location details"
p.level = 1

p = tf.add_paragraph()
p.text = "• Includes metrics: Sales, Profit, Quantity, Discount"
p.level = 1

p = tf.add_paragraph()
p.text = "• Time period: Multiple years of sales data"
p.level = 1

# Slide 4: Key Performance Indicators (KPIs)
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Key Performance Indicators (KPIs)"
tf = body_shape.text_frame
tf.text = "Core metrics chosen for business analysis:"

p = tf.add_paragraph()
p.text = "• Sales: Total revenue generated"
p.level = 1

p = tf.add_paragraph()
p.text = "• Profit: Net profit after costs and discounts"
p.level = 1

p = tf.add_paragraph()
p.text = "• Growth: Profit margin percentage"
p.level = 1

# Slide 5: Dashboard Features
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Dashboard Features"
tf = body_shape.text_frame
tf.text = "Interactive elements for enhanced user experience:"

p = tf.add_paragraph()
p.text = "• Slicers/Filters: Region, Category, Year range selection"
p.level = 1

p = tf.add_paragraph()
p.text = "• Time-series Analysis: Sales trends over time with line charts"
p.level = 1

p = tf.add_paragraph()
p.text = "• Cards for Totals: KPI summary cards with key metrics"
p.level = 1

p = tf.add_paragraph()
p.text = "• Consistent Color Theme: Blue-based professional color scheme"
p.level = 1

p = tf.add_paragraph()
p.text = "• Navigation Menu: Sidebar with filter controls"
p.level = 1

# Slide 6: Visualizations
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Visualizations Included"
tf = body_shape.text_frame
tf.text = "Charts and graphs for data analysis:"

p = tf.add_paragraph()
p.text = "• Bar Charts: Sales by Category, Top Products"
p.level = 1

p = tf.add_paragraph()
p.text = "• Pie Chart: Sales distribution by Region"
p.level = 1

p = tf.add_paragraph()
p.text = "• Line Chart: Sales trend over time"
p.level = 1

p = tf.add_paragraph()
p.text = "• Horizontal Bar Chart: Profit by Sub-Category"
p.level = 1

# Slide 7: Technical Implementation
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Technical Implementation"
tf = body_shape.text_frame
tf.text = "Built using Streamlit framework:"

p = tf.add_paragraph()
p.text = "• Python-based interactive web application"
p.level = 1

p = tf.add_paragraph()
p.text = "• Libraries: Pandas, Plotly, Streamlit"
p.level = 1

p = tf.add_paragraph()
p.text = "• Responsive design with sidebar navigation"
p.level = 1

p = tf.add_paragraph()
p.text = "• Real-time filtering and data updates"
p.level = 1

# Slide 8: Learning Outcomes
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Learning Outcomes"
tf = body_shape.text_frame
tf.text = "Skills developed through this project:"

p = tf.add_paragraph()
p.text = "• Dashboard design principles for business intelligence"
p.level = 1

p = tf.add_paragraph()
p.text = "• Data visualization techniques using modern tools"
p.level = 1

p = tf.add_paragraph()
p.text = "• Interactive filtering and user experience design"
p.level = 1

p = tf.add_paragraph()
p.text = "• KPI selection and business metric analysis"
p.level = 1

p = tf.add_paragraph()
p.text = "• Time-series analysis and trend identification"
p.level = 1

# Slide 9: Conclusion
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Conclusion"
tf = body_shape.text_frame
tf.text = "The interactive dashboard provides business stakeholders with:"

p = tf.add_paragraph()
p.text = "• Comprehensive view of sales performance"
p.level = 1

p = tf.add_paragraph()
p.text = "• Real-time filtering capabilities for detailed analysis"
p.level = 1

p = tf.add_paragraph()
p.text = "• Visual insights to support data-driven decisions"
p.level = 1

p = tf.add_paragraph()
p.text = "• Professional presentation with consistent theming"
p.level = 1

# Save presentation
prs.save('dashboard_summary.pptx')
print("PPT summary created successfully!")
